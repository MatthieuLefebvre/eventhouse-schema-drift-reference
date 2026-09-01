from __future__ import annotations

import argparse
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "presentation" / "eventhouse-schema-drift-reference.pptx"
ARCHITECTURE_IMAGE = ROOT / "docs" / "images" / "target-architecture.png"

BACKGROUND = RGBColor(14, 22, 29)
INK = RGBColor(8, 15, 20)
TEXT = RGBColor(235, 242, 244)
MUTED = RGBColor(170, 190, 198)
PANEL = RGBColor(27, 41, 50)
PANEL_BORDER = RGBColor(60, 82, 92)
BLUE = RGBColor(0, 120, 212)
TEAL = RGBColor(0, 133, 119)
AMBER = RGBColor(230, 126, 34)
RED = RGBColor(196, 43, 28)
WHITE = RGBColor(255, 255, 255)

SLIDES = [
    {"title": "Handling schema drift in Eventhouse", "bullets": ["A reusable pattern for JSON that changes over time", "Accept new data without making downstream schemas unpredictable"]},
    {"title": "Schema drift is more than a new field", "bullets": ["A producer adds a measurement", "A number starts arriving as text", "A property disappears from some messages", "A nested object changes shape", "An array contains a different number of items"], "visual": "drift"},
    {"title": "Flexible input, stable output", "bullets": ["Producers need room to evolve their messages", "Reports and APIs need columns they can rely on", "Rejecting every change loses useful data", "Accepting every key as a column creates an unstable contract", "The design has to serve both sides"], "visual": "contract"},
    {"title": "The five principles", "bullets": ["Keep the original message for investigation and replay", "Separate the input shape from the consumer contract", "Preserve unfamiliar values in residual JSON", "Detect drift without stopping known-field processing", "Promote a field only after review"], "visual": "principles"},
    {"title": "How the pieces fit together", "bullets": ["RawTelemetry is the recovery point", "Stored functions create fixed typed rows", "Update policies run those functions on arrival", "A separate path records drift for review", "Arrays are written as child rows"], "visual": "architecture"},
    {"title": "A new field arrives", "bullets": ["Controller 01 is still reporting its usual measurements", "Firmware version 2 adds serviceCountdownHours: 120", "There is no matching target column yet", "The value must survive without holding up the rest of the row"], "code": '"sourceType":"controller",\n"schemaVersion":2,\n"telemetry": {\n  "controllerStatus":"running",\n  "engineHours":1251.0,\n  "fuelConsumption":8.2,\n  "serviceCountdownHours":120\n}'},
    {"title": "First, keep the whole payload", "bullets": ["A few envelope values are mapped to columns", "DropMappedFields leaves everything else in RawRecord", "The landing table does not change when the payload does"], "code": 'SourceType     controller\nSchemaVersion  2\nRawRecord      {\n "payload":{"telemetry":{\n  ...,\n  "serviceCountdownHours":120\n }}}'},
    {"title": "Then write the fields we know", "bullets": ["The transform casts each approved value explicitly", "bag_remove_keys leaves the new property in ResidualTelemetry", "Reports continue to see the same columns as before"], "code": 'ControllerStatus  running\nEngineHours       1251.0\nFuelConsumption   8.2\nResidualTelemetry {\n "serviceCountdownHours":120\n}'},
    {"title": "Update policies do the routine work", "bullets": ["One policy writes the controller row", "Another checks the same message for unknown fields", "There is no notebook schedule or export watermark in between"], "code": '{"Source":"RawTelemetry",\n "Query":"TransformControllerTelemetry()",\n "IsTransactional":false}\n\nRaw row  →  typed row\n         →  drift observation'},
    {"title": "The new field goes onto the review list", "bullets": ["bag_keys finds serviceCountdownHours", "gettype records long; the sample value is 120", "Normal controller processing carries on while the team reviews it"], "code": 'SourceType    controller\nFieldPath     serviceCountdownHours\nObservedType  long\nSampleValue   120\n\nbag_keys + mv-expand\n+ set_has_element + gettype'},
    {"title": "Arrays are handled as rows", "bullets": ["This cooling-unit message has two zones", "mv-expand writes two child rows and keeps the zone IDs", "A message with an empty zones array writes no child rows"], "code": 'Device      Zone  Mode     Return  Setpoint\ncooling-01  1     cool     2.8     2.0\ncooling-01  2     defrost  5.1     4.0\n\n| mv-expand Zone=Zones'},
    {"title": "Promotion is a normal code change", "bullets": ["The team agrees the name, meaning, unit, and type", "The PR adds the column and updates the stored function", "After a schema check, only the agreed history is replayed"], "code": 'BEFORE\nResidual {"serviceCountdownHours":120}\n\nAFTER\nServiceCountdownHours  120.0\nResidual               {}\n\n.alter-merge + .set-or-append'},
    {"title": "An alert starts the conversation", "bullets": ["Do not alert for every message; aggregate first", "Here, 10 observations in 15 minutes is enough to open a ticket", "The ticket carries the type, sample, asset count, and time range"], "code": 'controller\nserviceCountdownHours\nType       long\nSample     120\nAssets     47\nEvents     8,212\n\n→ Teams  → DATA-1842'},
    {"title": "Who looks at the ticket?", "bullets": ["Operations checks that ingestion is healthy", "The source owner confirms the firmware change", "The domain owner explains what the field means", "Engineering profiles the data and prepares the change", "Consumer and change owners sign off before deployment"], "visual": "review"},
    {"title": "What the on-call engineer sees", "bullets": ["Is data arriving, and are typed tables keeping up?", "Which new fields appeared after the latest release?", "Are known fields failing type conversion?", "How large is the residual backlog?", "Are arrays producing an unexpected number of rows?"], "visual": "dashboard"},
    {"title": "Not every field becomes a column", "bullets": ["Promote it when the meaning is stable and consumers need it", "Keep it in residual JSON when it is sparse or diagnostic", "Reject or quarantine accidental, malformed, or sensitive data", "Record the decision so the same ticket does not return tomorrow"], "visual": "decision"},
    {"title": "Where this can simplify an existing design", "bullets": ["Some teams currently flatten JSON in Spark, pipelines, or application code", "Keep that processing when it adds capabilities KQL does not provide", "Move routine parsing and drift observation closer to ingestion when it reduces data movement", "Compare both paths with representative volume before cutting over"], "visual": "simplify"},
    {"title": "A sensible way to start", "bullets": ["Shadow one source type before changing the current path", "Measure ingestion cost, latency, array growth, and completeness", "Exercise failure and replay while the blast radius is small", "Change the existing process only after the numbers and outputs agree"], "visual": "adopt"},
]


def load_font(size: int):
    candidates = (
        "arial.ttf",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        "/usr/share/fonts/truetype/liberation2/LiberationSans-Regular.ttf",
    )
    for candidate in candidates:
        try:
            return ImageFont.truetype(candidate, size)
        except OSError:
            continue
    return ImageFont.load_default(size=size)


def add_textbox(slide, x, y, width, height, text, size=20, color=TEXT, bold=False):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(height))
    frame = shape.text_frame
    frame.clear()
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.font.name = "Aptos"
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color
    return shape


def add_header(slide, title, number):
    add_textbox(slide, 0.7, 0.45, 11.8, 0.6, title, 28, TEXT, True)
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(1.16), Inches(1.15), Inches(0.07))
    accent.fill.solid()
    accent.fill.fore_color.rgb = BLUE
    accent.line.fill.background()
    add_textbox(slide, 12.2, 0.5, 0.5, 0.4, f"{number:02}", 11, BLUE, True)


def add_bullets(slide, bullets):
    shape = slide.shapes.add_textbox(Inches(0.95), Inches(1.55), Inches(7.0), Inches(4.95))
    frame = shape.text_frame
    frame.clear()
    for index, bullet in enumerate(bullets):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.font.name = "Aptos"
        paragraph.font.size = Pt(20 if len(bullets) > 5 else 22)
        paragraph.font.color.rgb = TEXT
        paragraph.space_after = Pt(18)
        paragraph.text = f"•  {bullet}"


def add_code(slide, code):
    panel = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.25), Inches(1.55), Inches(4.45), Inches(4.8))
    panel.fill.solid()
    panel.fill.fore_color.rgb = INK
    panel.line.color.rgb = PANEL_BORDER
    frame = panel.text_frame
    frame.clear()
    frame.margin_left = Inches(0.25)
    frame.margin_right = Inches(0.2)
    frame.margin_top = Inches(0.25)
    paragraph = frame.paragraphs[0]
    paragraph.text = code
    paragraph.font.name = "Cascadia Mono"
    paragraph.font.size = Pt(14)
    paragraph.font.color.rgb = RGBColor(220, 238, 241)


def add_visual(slide, visual):
    panel = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.55), Inches(1.5), Inches(4.1), Inches(4.95))
    panel.fill.solid()
    panel.fill.fore_color.rgb = PANEL
    panel.line.color.rgb = PANEL_BORDER
    labels = {
        "spark": ("EXPORT → SPARK", RED),
        "cost": ("MOVE + INFER + MERGE", AMBER),
        "drift": ("ADD | REMOVE | RETYPE", AMBER),
        "contract": ("FLEXIBLE → STABLE", BLUE),
        "principles": ("KEEP → DETECT → REVIEW", TEAL),
        "journey": ("BUILD → PROVE → PROMOTE", BLUE),
        "proof": ("DRIFT SURVIVES", TEAL),
        "review": ("PROFILE → APPROVE", AMBER),
        "dashboard": ("OPERATIONS VIEW", TEAL),
        "decision": ("3 DECISIONS", BLUE),
        "simplify": ("MOVE ONLY WHAT FITS", AMBER),
        "adopt": ("SHADOW → BENCHMARK", BLUE),
    }
    word, color = labels.get(visual, ("EVENTHOUSE", BLUE))
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.15), Inches(3.1), Inches(2.9), Inches(1.15))
    badge.fill.solid()
    badge.fill.fore_color.rgb = color
    badge.line.fill.background()
    frame = badge.text_frame
    frame.clear()
    paragraph = frame.paragraphs[0]
    paragraph.text = word
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.font.name = "Aptos Display"
    paragraph.font.size = Pt(20)
    paragraph.font.bold = True
    paragraph.font.color.rgb = WHITE


def add_architecture(slide):
    slide.shapes.add_picture(str(ARCHITECTURE_IMAGE), Inches(8.15), Inches(1.45), width=Inches(4.65), height=Inches(3.95))
    add_textbox(slide, 8.25, 5.65, 4.35, 0.55, "No routine telemetry export", 17, TEAL, True)


def build_architecture_image():
    ARCHITECTURE_IMAGE.parent.mkdir(parents=True, exist_ok=True)
    image = Image.new("RGB", (1600, 900), "#0e161d")
    draw = ImageDraw.Draw(image)
    font = load_font(31)
    small = load_font(24)
    boxes = [
        ((70, 350, 300, 500), "Event Hub", "#0078d4"),
        ((370, 350, 640, 500), "RawTelemetry", "#182530"),
        ((720, 100, 1050, 250), "Typed policies", "#008577"),
        ((720, 375, 1050, 525), "Zone policy", "#008577"),
        ((720, 650, 1050, 800), "Drift policy", "#e67e22"),
        ((1130, 100, 1510, 250), "Typed tables", "#0078d4"),
        ((1130, 375, 1510, 525), "Child rows", "#0078d4"),
        ((1130, 650, 1510, 800), "Drift observations", "#e67e22"),
    ]
    for (left, top, right, bottom), label, color in boxes:
        draw.rounded_rectangle((left, top, right, bottom), radius=12, fill=color)
        text_box = draw.textbbox((0, 0), label, font=font)
        draw.text(((left + right - text_box[2]) / 2, (top + bottom - text_box[3]) / 2), label, font=font, fill="white")
    arrows = [((300, 425), (370, 425)), ((640, 425), (720, 175)), ((640, 425), (720, 450)), ((640, 425), (720, 725)), ((1050, 175), (1130, 175)), ((1050, 450), (1130, 450)), ((1050, 725), (1130, 725))]
    for start, end in arrows:
        draw.line((start, end), fill="#9bb0b8", width=6)
        draw.polygon([(end[0], end[1]), (end[0] - 18, end[1] - 10), (end[0] - 18, end[1] + 10)], fill="#9bb0b8")
    draw.text((70, 65), "Eventhouse-native schema drift", font=font, fill="#ebf2f4")
    draw.text((70, 115), "Fixed target schemas • residual JSON • reviewed promotion", font=small, fill="#aabeC6")
    image.save(ARCHITECTURE_IMAGE)


def build_presentation(output: Path = OUTPUT):
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    for number, content in enumerate(SLIDES, 1):
        title = content["title"]
        bullets = content["bullets"]
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        background = slide.background.fill
        background.solid()
        background.fore_color.rgb = BACKGROUND
        if number == 1:
            stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.22), presentation.slide_height)
            stripe.fill.solid()
            stripe.fill.fore_color.rgb = BLUE
            stripe.line.fill.background()
            add_textbox(slide, 0.9, 1.5, 11.5, 1.3, title, 38, TEXT, True)
            add_textbox(slide, 0.95, 3.15, 10.8, 1.3, "\n".join(bullets), 21, MUTED)
            add_textbox(slide, 0.95, 6.45, 5.0, 0.4, "PUBLIC REFERENCE IMPLEMENTATION", 11, BLUE, True)
        else:
            add_header(slide, title, number)
            add_bullets(slide, bullets)
            if "code" in content:
                add_code(slide, content["code"])
            elif content.get("visual") == "architecture":
                add_architecture(slide)
            else:
                add_visual(slide, content.get("visual"))
    output.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(output)


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Build the customer demo presentation.")
    parser.add_argument("--output", type=Path, default=OUTPUT)
    arguments = parser.parse_args()
    build_architecture_image()
    build_presentation(arguments.output)
    print(f"wrote {arguments.output}")
    print(f"wrote {ARCHITECTURE_IMAGE.relative_to(ROOT)}")
